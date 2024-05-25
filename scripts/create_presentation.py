import os
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import random
from tqdm import tqdm

PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
current_dir = os.getcwd()


def add_flashcard(question, answer, output_slide, slide_layout, conf):
    new_slide = output_slide.slides.add_slide(slide_layout)
    fill = new_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(223, 223, 223)

    flash_card = conf["flash_card"]

    bg = random.choice(conf["bgs"])
    bg_img_path = os.path.join(PROJECT_ROOT, bg["img_path"])
    new_slide.shapes.add_picture(
        bg_img_path, bg["left"], bg["top"], bg["width"], bg["height"]
    )

    flash_card_img_path = os.path.join(PROJECT_ROOT, flash_card["img_path"])
    new_slide.shapes.add_picture(
        flash_card_img_path,
        flash_card["left"],
        flash_card["top"],
        flash_card["width"],
        flash_card["height"],
    )

    for name, text_shape in conf["texts"].items():
        new_textbox = new_slide.shapes.add_textbox(
            text_shape["left"],
            text_shape["top"],
            text_shape["width"],
            text_shape["height"],
        )
        new_text_frame = new_textbox.text_frame
        new_paragraph = new_text_frame.paragraphs[0]

        # add text
        if name == "at":
            if isinstance(answer, str):
                new_paragraph.text = answer
                new_paragraph.font.color.rgb = RGBColor(
                    text_shape["color"][0],
                    text_shape["color"][1],
                    text_shape["color"][2],
                )
                if len(answer) > 217:
                    new_paragraph.font.size = Pt(33)
            else:
                options = answer["options"]
                correct_answer = answer["correct_answer"]
                new_paragraph.text = ", ".join(options).capitalize()
                new_paragraph.font.color.rgb = RGBColor(255, 0, 0)

                new_paragraph2 = new_text_frame.add_paragraph()
                new_paragraph2.text = f"Correct answer: {correct_answer}"
                new_paragraph2.font.color.rgb = RGBColor(0, 255, 0)
                new_paragraph2.alignment = text_shape["alignment"]
                if len("".join(options)) + len(correct_answer) > 217:
                    new_paragraph.font.size = Pt(33)

        elif name == "qt":
            new_paragraph.text = question
            new_paragraph.font.color.rgb = RGBColor(
                text_shape["color"][0], text_shape["color"][1], text_shape["color"][2]
            )

            if len(question) > 158:
                new_paragraph.font.size = Pt(36)

        new_text_frame.fit_text(
            font_family=text_shape["font"],
            max_size=text_shape["max_size"],
            font_file=os.path.join(PROJECT_ROOT, text_shape["font_file"]),
        )
        new_text_frame.word_wrap = text_shape["word_wrap"]
        new_paragraph.alignment = text_shape["alignment"]
        new_text_frame.vertical_anchor = text_shape["vertical_anchor"]
        new_text_frame.auto_size = text_shape["auto_size"]


def add_header_slide(title, conf, slide_layout, output_slide):
    new_slide = output_slide.slides.add_slide(slide_layout)

    header = conf["header"]

    header_img_path = os.path.join(PROJECT_ROOT, header["img_path"])
    new_slide.shapes.add_picture(
        header_img_path,
        header["left"],
        header["top"],
        header["width"],
        header["height"],
    )

    title_ = header["title"]

    new_textbox = new_slide.shapes.add_textbox(
        title_["left"], title_["top"], title_["width"], title_["height"]
    )
    new_text_frame = new_textbox.text_frame
    new_paragraph = new_text_frame.paragraphs[0]
    new_paragraph.alignment = title_["alignment"]
    new_run = new_paragraph.add_run()
    new_run.text = title
    new_run.font.size = title_["size"]
    new_run.font.color.rgb = RGBColor(
        title_["color"][0], title_["color"][1], title_["color"][2]
    )
    new_run.font.name = title_["font"]
    new_text_frame.word_wrap = True


def add_footer_slide(conf, slide_layout, output_slide):
    new_slide = output_slide.slides.add_slide(slide_layout)

    footer = conf["footer"]

    footer_img_path = os.path.join(PROJECT_ROOT, footer["img_path"])
    new_slide.shapes.add_picture(
        footer_img_path,
        footer["left"],
        footer["top"],
        footer["width"],
        footer["height"],
    )


def create(output_path, output_pptx_name, json_path, title="Flashcards"):
    """
    Creates a PowerPoint presentation with flashcards based on the provided JSON data.

    Args:
        output_path (str): The path where the output presentation will be saved.
        output_pptx_name (str): The name of the output PowerPoint presentation file.
        json_path (str or list): The path(s) to the JSON file(s) containing flashcard data.
            If a single string is provided, it is treated as the path to a single JSON file.
            If a list of strings is provided, each string is treated as the path to a JSON file.
        conf_path (str, optional): The path to the configuration JSON file. Defaults to "static/conf.json".
        title (str, optional): The title of the presentation. Defaults to "Flashcards".

    Returns:
        None
    """
    conf_path = os.path.join(PROJECT_ROOT,"static/conf.json")

    with open(conf_path, "r") as f:
        conf = json.load(f)
    output_pptx = Presentation()
    output_pptx.slide_height, output_pptx.slide_width = (
        conf["slide_height"],
        conf["slide_width"],
    )
    slide_layout = output_pptx.slide_layouts[5]  # Assuming layout index 5 is "Blank"

    add_header_slide(title, conf, slide_layout, output_pptx)

    json_files = []
    if not isinstance(json_path, list):
        if os.path.isfile(json_path):
            json_files.append(json_path)
        elif os.path.isdir(json_path):
            for file in os.listdir(json_path):
                if file.endswith(".json"):
                    json_files.append(os.path.join(json_path, file))
    else:
        json_files = json_path

    for json_file in tqdm(json_files, desc="Processing JSON files"):
        with open(json_file, "r") as f:
            data = json.load(f)
            flashcards = data["flashcards"]

        for flashcard in flashcards:
            if "question" in flashcard and "answer" in flashcard:
                add_flashcard(
                    flashcard["question"],
                    flashcard["answer"],
                    output_pptx,
                    slide_layout,
                    conf,
                )
            elif (
                "question" in flashcard
                and "options" in flashcard
                and "correct_answer" in flashcard
            ):
                add_flashcard(
                    flashcard["question"],
                    {
                        "options": flashcard["options"],
                        "correct_answer": flashcard["correct_answer"],
                    },
                    output_pptx,
                    slide_layout,
                    conf,
                )
            else:
                print(
                    "Each dictionary in the JSON file must contain 'question' and 'answer' keys."
                )

    add_footer_slide(conf, slide_layout, output_pptx)

    if not output_path:
        output_path = "."

    output = os.path.join(current_dir, output_path)
    if not os.path.exists(output):
        os.makedirs(output)

    output_file_path = os.path.join(output, output_pptx_name)
    if os.path.exists(output_file_path):
        os.remove(output_file_path)

    output_pptx.save(output_file_path)
    print(f"Flashcards created: {output_file_path}")
