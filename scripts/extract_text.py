#!/usr/bin/env python3
import os
import re
from typing import List
from pptx import Presentation
from pypdf import PdfReader
import docx2txt
PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))

# Function to clean up extracted text
def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    text = "".join([char if ord(char) < 128 else " " for char in text])
    return text.strip()

# Function to process and split text into chunks
def process_text(extracted_text: List[str], character_limit: int) -> List[str]:
    cleaned_text_chunks: List[str] = []
    for text in extracted_text:
        cleaned_text = clean_text(text)
        chunks = [
            cleaned_text[i : i + character_limit]
            for i in range(0, len(cleaned_text), character_limit)
        ]
        cleaned_text_chunks.extend(chunks)
    return cleaned_text_chunks

# Extraction functions for different file types
def extract_text_from_txt(txt_file_path: str) -> str:
    with open(txt_file_path, "r", encoding="utf-8") as file:
        return file.read()

def extract_text_from_docx(docx_file_path: str) -> str:
    return docx2txt.process(docx_file_path)

def extract_text_from_pdf(pdf_file_path: str) -> List[str]:
    reader = PdfReader(pdf_file_path)
    return [page.extract_text() for page in reader.pages]

def extract_text_from_pptx(pptx_file_path: str) -> List[str]:
    presentation = Presentation(pptx_file_path)
    extracted_text = []
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text.strip() + "\n"
        extracted_text.append(slide_text)
    return extracted_text


# Main function to process files and save the output
def extract(input_path: str,  character_limit: int= 5000):
    """
    Extracts text from files in the specified input path and saves the extracted text
    into multiple output files based on the character limit.

    Args:
        input_path (str): The path to the input file or directory.
        output_path (str): The path to the output directory where the extracted text files will be saved.
        character_limit (int): The maximum number of characters allowed in each output file.

    Raises:
        ValueError: If the input path is invalid.

    Returns:
        None
    """
    output_path = os.path.join(PROJECT_ROOT, "extracted_data")

    if not os.path.exists(output_path):
        os.makedirs(output_path)
    
    extracted_data = []
    if os.path.isfile(input_path):
        if input_path.endswith(".pptx"):
            extracted_data.extend(extract_text_from_pptx(input_path))
        elif input_path.endswith(".pdf"):
            extracted_data.extend(extract_text_from_pdf(input_path))
        elif input_path.endswith(".txt"):
            extracted_data.append(extract_text_from_txt(input_path))
        elif input_path.endswith(".docx"):
            extracted_data.append(extract_text_from_docx(input_path))
    elif os.path.isdir(input_path):
        for file in os.listdir(input_path):
            file_path = os.path.join(input_path, file)
            if file.endswith(".pptx"):
                extracted_data.extend(extract_text_from_pptx(file_path))
            elif file.endswith(".pdf"):
                extracted_data.extend(extract_text_from_pdf(file_path))
            elif file.endswith(".txt"):
                extracted_data.append(extract_text_from_txt(file_path))
            elif file.endswith(".docx"):
                extracted_data.append(extract_text_from_docx(file_path))
    else:
        raise ValueError("Invalid input path.")
    
    cleaned_text_chunks = process_text(extracted_data, character_limit)

    buffer = ""
    file_count = 0
    for chunk in cleaned_text_chunks:
        if len(buffer + chunk) > character_limit:
            with open(
                os.path.join(output_path, f"output_{file_count}.txt"),
                "w",
                encoding="utf-8",
            ) as output_file:
                output_file.write(buffer)
            buffer = ""
            file_count += 1
        buffer += chunk
    
    if buffer:
        with open(
            os.path.join(output_path, f"output_{file_count}.txt"),
            "w",
            encoding="utf-8",
        ) as output_file:
            output_file.write(buffer)

    print(f"Output files saved to {output_path}.")


