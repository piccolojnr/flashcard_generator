from pptx import Presentation
import json
from pptx.dml.color import RGBColor
from pptx.util import Pt
import os
import random

script_dir = os.path.dirname(os.path.abspath(__file__))
script_dir = os.path.dirname(script_dir)
current_dir = os.getcwd()


def add_flashcard(question, answer, output_slide, slide_layout, conf):
    new_slide = output_slide.slides.add_slide(slide_layout)
    fill = new_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(223, 223, 223)

    flash_card = conf["flash_card"]

    bg = random.choice(conf["bgs"])
    bg_img_path = os.path.join(script_dir, bg["img_path"])
    new_slide.shapes.add_picture(
        bg_img_path, bg["left"], bg["top"], bg["width"], bg["height"]
    )

    flash_card_img_path = os.path.join(script_dir, flash_card["img_path"])
    new_slide.shapes.add_picture(
        flash_card_img_path,
        flash_card["left"],
        flash_card["top"],
        flash_card["width"],
        flash_card["height"],
    )

    for name, text_shape in conf["texts"].items():
        new_textbox = new_slide.shapes.add_textbox(
            text_shape["left"],
            text_shape["top"],
            text_shape["width"],
            text_shape["height"],
        )
        new_text_frame = new_textbox.text_frame
        new_paragraph = new_text_frame.paragraphs[0]

        # add text
        if name == "at":
            if type(answer) == str:
                new_paragraph.text = answer
                new_paragraph.font.color.rgb = RGBColor(
                    text_shape["color"][0],
                    text_shape["color"][1],
                    text_shape["color"][2],
                )
                if len(answer) > 217:
                    new_paragraph.font.size = Pt(33)
            else:
                options = answer["options"]
                correct_answer = answer["correct_answer"]
                new_paragraph.text = ", ".join(options).capitalize()
                new_paragraph.font.color.rgb = RGBColor(255, 0, 0)

                new_paragraph2 = new_text_frame.add_paragraph()
                new_paragraph2.text = f"Correct answer: {correct_answer}"
                new_paragraph2.font.color.rgb = RGBColor(0, 255, 0)
                new_paragraph2.alignment = text_shape["alignment"]
                if len("".join(options)) + len(correct_answer) > 217:
                    new_paragraph.font.size = Pt(33)

        elif name == "qt":
            new_paragraph.text = question
            new_paragraph.font.color.rgb = RGBColor(
                text_shape["color"][0], text_shape["color"][1], text_shape["color"][2]
            )

            if len(question) > 158:
                new_paragraph.font.size = Pt(36)

        new_text_frame.fit_text(
            font_family=text_shape["font"],
            max_size=text_shape["max_size"],
            font_file=os.path.join(script_dir, text_shape["font_file"]),
        )
        new_text_frame.word_wrap = text_shape["word_wrap"]
        new_paragraph.alignment = text_shape["alignment"]
        new_text_frame.vertical_anchor = text_shape["vertical_anchor"]
        new_text_frame.auto_size = text_shape["auto_size"]


def add_header_slide(title, conf, slide_layout, output_slide):
    new_slide = output_slide.slides.add_slide(slide_layout)

    header = conf["header"]

    header_img_path = os.path.join(script_dir, header["img_path"])
    new_slide.shapes.add_picture(
        header_img_path,
        header["left"],
        header["top"],
        header["width"],
        header["height"],
    )

    title_ = header["title"]

    new_textbox = new_slide.shapes.add_textbox(
        title_["left"], title_["top"], title_["width"], title_["height"]
    )
    new_text_frame = new_textbox.text_frame
    new_paragraph = new_text_frame.paragraphs[0]
    new_paragraph.alignment = title_["alignment"]
    new_run = new_paragraph.add_run()
    new_run.text = title
    new_run.font.size = title_["size"]
    new_run.font.color.rgb = RGBColor(
        title_["color"][0], title_["color"][1], title_["color"][2]
    )
    new_run.font.name = title_["font"]
    new_text_frame.word_wrap = True


def add_footer_slide(conf, slide_layout, output_slide):
    new_slide = output_slide.slides.add_slide(slide_layout)

    footer = conf["footer"]

    footer_img_path = os.path.join(script_dir, footer["img_path"])
    new_slide.shapes.add_picture(
        footer_img_path,
        footer["left"],
        footer["top"],
        footer["width"],
        footer["height"],
    )


def main(
    output_path,
    output_pptx_name,
    json_path,
    conf_path="static/conf.json",
    title="Flashcards",
):

    with open(conf_path, "r") as f:
        conf = json.load(f)

    output_pptx = Presentation()
    output_pptx.slide_height, output_pptx.slide_width = (
        conf["slide_height"],
        conf["slide_width"],
    )
    slide_layout = output_pptx.slide_layouts.get_by_name("Blank")

    add_header_slide(title, conf, slide_layout, output_pptx)

    json_files = []

    if os.path.isfile(json_path):
        json_files.append(json_path)
    elif os.path.isdir(json_path):
        for file in os.listdir(json_path):
            if file.endswith(".json"):
                json_files.append(os.path.join(json_path, file))

    for json_file in json_files:
        with open(json_file, "r") as f:
            flashcards = json.load(f)

        for flashcard in flashcards:
            if "question" in flashcard and "answer" in flashcard:
                add_flashcard(
                    flashcard["question"],
                    flashcard["answer"],
                    output_pptx,
                    slide_layout,
                    conf,
                )
            elif (
                "question" in flashcard
                and "options" in flashcard
                and "correct_answer" in flashcard
            ):
                add_flashcard(
                    flashcard["question"],
                    {
                        "options": flashcard["options"],
                        "correct_answer": flashcard["correct_answer"],
                    },
                    output_pptx,
                    slide_layout,
                    conf,
                )
            else:

                print(
                    "Each dictionary in the JSON file must contain 'question' and 'answer' keys."
                )

    add_footer_slide(conf, slide_layout, output_pptx)

    if not output_path:
        output_path = "."

    output = os.path.join(current_dir, output_path)
    if not os.path.exists(output):
        os.makedirs(output)

    if os.path.exists(os.path.join(output, output_pptx_name)):
        os.remove(os.path.join(output, output_pptx_name))

    output_pptx.save(os.path.join(output, output_pptx_name))
    print(f"Flashcards created: {output_pptx_name}")


def create_flashcards(output_path, output_name, title, conf_path="static/conf.json"):
    try:
        conf_path = os.path.join(script_dir, conf_path)

        if not os.path.exists(conf_path):
            raise FileNotFoundError(f"Configuration file not found: {conf_path}")

        if len(os.listdir("./flashcard_data")) == 0:
            print("No flashcards found in ./flashcard_data")
            return

        main(
            output_path,
            output_name + ".pptx",
            "./flashcard_data",
            conf_path,
            title,
        )

        print(f"Flashcards created: {output_name}")
        return os.path.join(output_path, output_name + ".pptx")
    except Exception as e:
        import traceback

        traceback.print_exc()
        print(f"Error: {e}")
        exit()


def create_flashcards_from_args(args):
    try:
        conf_path = os.path.join(script_dir, args.config)
        if not os.path.exists(conf_path):
            raise FileNotFoundError(f"Configuration file not found: {conf_path}")

        main(None, args.output + ".pptx", args.input, conf_path, args.title)
    except Exception as e:
        import traceback

        traceback.print_exc()
        print(f"Error: {e}")
        exit()
